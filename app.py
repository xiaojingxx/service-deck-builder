import os
import io
import base64
import tempfile
import subprocess
import shutil
from io import BytesIO
from shutil import which

import fitz  # PyMuPDF
import gspread
import psutil
import re
import streamlit as st
from streamlit_ace import st_ace
from google.oauth2.service_account import Credentials
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt
from PIL import Image
from docx import Document


# =========================================================
# CONFIG
# =========================================================
SOFFICE_PATH = os.environ.get("SOFFICE_PATH", "soffice")

SHEET_KEY = st.secrets["SHEET_KEY"]
WORKSHEET_NAME = st.secrets["WORKSHEET_NAME"]

FIRST_LAYOUT_NAME = "TEMPLATE_FIRST"
REST_LAYOUT_NAME = "TEMPLATE_REST"
DIVIDER_LAYOUT_NAME = "SECTION_DIVIDER"


# =========================================================
# STREAMLIT PAGE
# =========================================================
st.set_page_config(page_title="Service Deck Builder", layout="wide")
st.title("Service Deck Builder")


# =========================================================
# GOOGLE SHEETS
# =========================================================
@st.cache_resource
def get_sheet():
    scopes = [
        "https://www.googleapis.com/auth/spreadsheets",
        "https://www.googleapis.com/auth/drive",
    ]
    credentials = Credentials.from_service_account_info(
        st.secrets["gcp_service_account"],
        scopes=scopes,
    )
    gc = gspread.authorize(credentials)
    return gc.open_by_key(SHEET_KEY).worksheet(WORKSHEET_NAME)


sheet = get_sheet()


@st.cache_data(show_spinner=False)
def get_all_records_cached():
    return sheet.get_all_records()


# =========================================================
# SESSION STATE
# =========================================================
DEFAULTS = {
    "setlist": [],
    "loaded_song": None,
    "editing_setlist_index": None,
    "pending_setlist_load": None,
    "reset_editor_pending": False,
    "setlist_selected_index": 0,
    "uploaded_templates": {},
    "selected_template_name": None,
    "editor_umh": "",
    "editor_title": "",
    "editor_text": "",
    "editor_override_lyrics_font_size": False,
    "editor_override_line_spacing": False,
    "editor_lyrics_font_size_pt": 32,
    "editor_line_spacing": 1.2,
    "auto_split_by_lines": False,
    "lines_per_slide": 4,
    "refresh_on_new_line": True,
    "editor_ace_key": 0,
    "last_editor_text": "",
    "last_current_song_signature": None,
    "editor_status_message": "",
    "current_preview_slide": 1,
    "preview_mode": "song",
    "current_song_preview_images": None,
    "current_song_preview_stats": None,
    "service_preview_images": None,
    "service_preview_stats": None,
    "service_preview_error": None,
    "service_song_start_slides": [],
    "ppt_data": None,
    "last_split_settings": None,
    "preserve_template_slides": True,
    "template_sections": [],
    "hidden_section_ids": [],
    "service_output_mode": "full",  # full | songs
    "selected_song_section_id": None,
    "setlist_selectbox_sidebar": 0,
    "pending_setlist_selectbox_index": None,
    "service_order_blocks": [],
    "song_store": {},
    "last_docx_section_mapping": [],
    "selected_song_id": None,
}

for key, value in DEFAULTS.items():
    if key not in st.session_state:
        st.session_state[key] = value


# =========================================================
# BASIC HELPERS
# =========================================================
def soffice_available() -> bool:
    if SOFFICE_PATH == "soffice":
        return which("soffice") is not None
    return os.path.exists(SOFFICE_PATH)


def clear_service_outputs():
    st.session_state["ppt_data"] = None
    st.session_state["service_preview_images"] = None
    st.session_state["service_preview_stats"] = None
    st.session_state["service_preview_error"] = None
    st.session_state["service_song_start_slides"] = []


def format_bytes(num_bytes: int) -> str:
    if num_bytes < 1024:
        return f"{num_bytes} B"
    if num_bytes < 1024**2:
        return f"{num_bytes / 1024:.1f} KB"
    if num_bytes < 1024**3:
        return f"{num_bytes / (1024**2):.2f} MB"
    return f"{num_bytes / (1024**3):.2f} GB"


def preview_stats(images):
    if not images:
        return {
            "count": 0,
            "total_bytes": 0,
            "avg_bytes": 0,
            "max_bytes": 0,
        }

    sizes = [len(img) for img in images]
    total = sum(sizes)

    return {
        "count": len(images),
        "total_bytes": total,
        "avg_bytes": total // len(images),
        "max_bytes": max(sizes),
    }


def get_runtime_resource_stats():
    proc = psutil.Process(os.getpid())
    mem = proc.memory_info()

    vm = psutil.virtual_memory()
    disk = shutil.disk_usage("/")

    return {
        "process_rss": mem.rss,
        "process_vms": mem.vms,
        "system_available_ram": vm.available,
        "system_total_ram": vm.total,
        "system_used_ram": vm.used,
        "disk_free": disk.free,
        "disk_total": disk.total,
        "disk_used": disk.used,
    }


def find_row_by_umh(umh_number: str):
    umh_number = str(umh_number).strip()
    for row in get_all_records_cached():
        if str(row.get("UMH Number", "")).strip() == umh_number:
            return row
    return None


def search_titles(keyword: str, limit: int = 20):
    keyword = keyword.lower().strip()
    matches = []
    if not keyword:
        return matches

    for row in get_all_records_cached():
        title = str(row.get("Title", "")).strip()
        if keyword in title.lower():
            matches.append(row)

    return matches[:limit]


def split_slides_manual(text: str) -> list[list[str]]:
    blocks = [block.strip() for block in text.split("\n\n") if block.strip()]
    slides = []

    for block in blocks:
        lines = [line.strip() for line in block.splitlines() if line.strip()]
        if lines:
            slides.append(lines)

    return slides


def split_slides_by_line_count(text: str, lines_per_slide: int = 4) -> list[list[str]]:
    if lines_per_slide < 1:
        lines_per_slide = 1

    verses = []
    current_verse = []

    for raw_line in text.splitlines():
        line = raw_line.strip()
        if line == "":
            if current_verse:
                verses.append(current_verse)
                current_verse = []
        else:
            current_verse.append(line)

    if current_verse:
        verses.append(current_verse)

    slides = []
    for verse in verses:
        for i in range(0, len(verse), lines_per_slide):
            chunk = verse[i:i + lines_per_slide]
            if chunk:
                slides.append(chunk)

    return slides


def get_current_slides(text: str) -> list[list[str]]:
    if st.session_state["auto_split_by_lines"]:
        return split_slides_by_line_count(
            text,
            lines_per_slide=st.session_state["lines_per_slide"],
        )
    return split_slides_manual(text)


def open_presentation_from_bytes(template_bytes: bytes):
    return Presentation(BytesIO(template_bytes))


def get_layout_by_name(prs, layout_name):
    for slide_master in prs.slide_masters:
        for slide_layout in slide_master.slide_layouts:
            if slide_layout.name.strip() == layout_name:
                return slide_layout
    return None


def get_body_placeholder(slide):
    placeholders = list(slide.placeholders)

    for ph in placeholders:
        name = getattr(ph, "name", "").lower()
        if "title" not in name and getattr(ph, "has_text_frame", False):
            return ph

    if len(placeholders) > 1:
        return placeholders[1]
    if placeholders:
        return placeholders[0]
    return None


def set_shape_text(shape, text, font_size_pt=None, line_spacing=None):
    if shape is None or not getattr(shape, "has_text_frame", False):
        return

    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True

    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER

        if line_spacing is not None:
            p.line_spacing = line_spacing

        run = p.add_run()
        run.text = line

        if font_size_pt is not None:
            run.font.size = Pt(font_size_pt)


def delete_all_slides(prs):
    while len(prs.slides) > 0:
        slide_id = prs.slides._sldIdLst[0]
        rId = slide_id.rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]


def delete_slide_by_index(prs, slide_index: int):
    slide_id = prs.slides._sldIdLst[slide_index]
    rId = slide_id.rId
    prs.part.drop_rel(rId)
    del prs.slides._sldIdLst[slide_index]

def apply_docx_heading_alias(heading: str) -> str:
    """
    Map common DOCX headings to template section names.
    This is a controlled alias layer (safe hardcoding).
    """
    simple = simplify_heading_text(heading)

    alias_map = {
        "closing hymn": "response",
        "closing song": "response",
        "hymn of response": "response",
        "response hymn": "response",

        "offertory": "tithe offering",
        "offering": "tithe offering",

        "announcements": "announcements",
        "welcome announcements": "announcements",

        "doxology": "doxology",
        "gloria patri": "gloria patri",
    }

    return alias_map.get(simple, heading)

def get_flat_song_index_by_song_id(setlist, song_id):
    if not song_id:
        return 0

    for i, song in enumerate(setlist):
        if song.get("song_id") == song_id:
            return i
    return 0


def restore_selected_index_from_song_id():
    setlist = st.session_state.get("setlist", [])
    restored_index = get_flat_song_index_by_song_id(
        setlist,
        st.session_state.get("selected_song_id"),
    )

    if setlist:
        restored_index = max(0, min(restored_index, len(setlist) - 1))
    else:
        restored_index = 0

    st.session_state["setlist_selected_index"] = restored_index
    st.session_state["pending_setlist_selectbox_index"] = restored_index

# =========================================================
# MOVE HELPERS
# =========================================================
def move_slide(prs, from_idx: int, to_idx: int):
    sldIdLst = prs.slides._sldIdLst

    if from_idx == to_idx:
        return

    slide_id = sldIdLst[from_idx]
    del sldIdLst[from_idx]

    if to_idx > from_idx:
        to_idx -= 1

    sldIdLst.insert(to_idx, slide_id)


def move_slide_block(prs, start_idx: int, end_idx: int, target_idx: int):
    sldIdLst = prs.slides._sldIdLst

    if start_idx > end_idx:
        start_idx, end_idx = end_idx, start_idx

    block_len = end_idx - start_idx + 1

    if target_idx >= start_idx and target_idx <= end_idx + 1:
        return

    block = [sldIdLst[i] for i in range(start_idx, end_idx + 1)]

    for _ in range(block_len):
        del sldIdLst[start_idx]

    if target_idx > end_idx:
        target_idx -= block_len

    for offset, slide_id in enumerate(block):
        sldIdLst.insert(target_idx + offset, slide_id)


# =========================================================
# TEMPLATE / SECTION HELPERS
# =========================================================
def get_slide_layout_name(slide):
    try:
        return slide.slide_layout.name.strip()
    except Exception:
        return ""


def is_divider_slide(slide):
    return get_slide_layout_name(slide) == DIVIDER_LAYOUT_NAME


def get_divider_title(slide, fallback="Untitled Section"):
    if slide.shapes.title and slide.shapes.title.text.strip():
        return slide.shapes.title.text.strip()

    texts = []
    for shape in slide.shapes:
        if hasattr(shape, "text"):
            text = shape.text.strip()
            if text:
                texts.append(text)

    return texts[0] if texts else fallback


def parse_template_sections(template_bytes: bytes):
    prs = open_presentation_from_bytes(template_bytes)

    sections = []
    current_section = None
    section_counter = 0

    for i, slide in enumerate(prs.slides):
        if is_divider_slide(slide):
            current_section = {
                "id": f"sec_{section_counter}",
                "title": get_divider_title(slide, fallback=f"Section {section_counter + 1}"),
                "divider_index": i,
                "content_slide_indices": [],
            }
            sections.append(current_section)
            section_counter += 1
        else:
            if current_section is not None:
                current_section["content_slide_indices"].append(i)

    return sections


def get_section_title_by_id(section_id):
    for sec in st.session_state.get("template_sections", []):
        if sec["id"] == section_id:
            return sec["title"]
    return None


def get_current_divider_positions(prs):
    positions = []
    for i, slide in enumerate(prs.slides):
        if is_divider_slide(slide):
            positions.append((i, get_divider_title(slide)))
    return positions


def canonicalize_section_label(s: str) -> str:
    s = str(s or "").strip().lower()
    s = re.sub(r"\s+", " ", s)
    return s



def find_section_insert_index(prs, section_title: str):
    divider_positions = get_current_divider_positions(prs)
    target_norm = canonicalize_section_label(section_title)

    match_index = None
    for idx, (_, title) in enumerate(divider_positions):
        if canonicalize_section_label(title) == target_norm:
            match_index = idx
            break

    if match_index is None:
        for idx, (_, title) in enumerate(divider_positions):
            title_norm = canonicalize_section_label(title)
            if target_norm and (target_norm in title_norm or title_norm in target_norm):
                match_index = idx
                break

    if match_index is None:
        return len(prs.slides)

    if match_index + 1 < len(divider_positions):
        return divider_positions[match_index + 1][0]

    return len(prs.slides)


def validate_template_bytes(template_bytes: bytes):
    prs = open_presentation_from_bytes(template_bytes)

    errors = []
    warnings = []

    first_layout = get_layout_by_name(prs, FIRST_LAYOUT_NAME)
    rest_layout = get_layout_by_name(prs, REST_LAYOUT_NAME)

    if first_layout is None:
        errors.append(f"Missing layout: {FIRST_LAYOUT_NAME}")
    if rest_layout is None:
        errors.append(f"Missing layout: {REST_LAYOUT_NAME}")

    if errors:
        return False, errors, warnings

    first_slide = prs.slides.add_slide(first_layout)
    rest_slide = prs.slides.add_slide(rest_layout)

    title_placeholder = first_slide.shapes.title
    first_body_placeholder = get_body_placeholder(first_slide)
    rest_body_placeholder = get_body_placeholder(rest_slide)

    if title_placeholder is None:
        errors.append(f"{FIRST_LAYOUT_NAME} is missing a title placeholder")
    if first_body_placeholder is None:
        errors.append(f"{FIRST_LAYOUT_NAME} is missing a body/lyrics placeholder")
    if rest_body_placeholder is None:
        errors.append(f"{REST_LAYOUT_NAME} is missing a body/lyrics placeholder")

    return len(errors) == 0, errors, warnings


# =========================================================
# SONG / SECTION HELPERS
# =========================================================
def build_editor_song_item(current_slides):
    return {
        "umh_number": st.session_state["editor_umh"].strip(),
        "title": st.session_state["editor_title"].strip(),
        "slides": current_slides,
        "lyrics_font_size_pt": (
            st.session_state["editor_lyrics_font_size_pt"]
            if st.session_state["editor_override_lyrics_font_size"]
            else None
        ),
        "line_spacing": (
            st.session_state["editor_line_spacing"]
            if st.session_state["editor_override_line_spacing"]
            else None
        ),
        "override_lyrics_font_size": st.session_state["editor_override_lyrics_font_size"],
        "override_line_spacing": st.session_state["editor_override_line_spacing"],
        "section_id": st.session_state.get("selected_song_section_id"),
    }


def build_current_song_signature(song_item, selected_template_name):
    return (
        song_item["umh_number"],
        song_item["title"],
        tuple(tuple(slide) for slide in song_item["slides"]),
        song_item.get("lyrics_font_size_pt"),
        song_item.get("line_spacing"),
        song_item.get("section_id"),
        selected_template_name,
        st.session_state["auto_split_by_lines"],
        st.session_state["lines_per_slide"],
    )


def get_ordered_songs_for_output(setlist):
    sections = st.session_state.get("template_sections", [])
    songs_by_section = {sec["id"]: [] for sec in sections}
    unassigned_songs = []

    for idx, song in enumerate(setlist):
        sec_id = song.get("section_id")
        if sec_id in songs_by_section:
            songs_by_section[sec_id].append((idx, song))
        else:
            unassigned_songs.append((idx, song))

    ordered = []
    for sec in sections:
        ordered.extend(songs_by_section.get(sec["id"], []))
    ordered.extend(unassigned_songs)
    return ordered


def group_songs_by_section_order(setlist):
    sections = st.session_state.get("template_sections", [])
    songs_by_section = {sec["id"]: [] for sec in sections}
    unassigned = []

    for idx, song in enumerate(setlist):
        sec_id = song.get("section_id")
        if sec_id in songs_by_section:
            songs_by_section[sec_id].append((idx, song))
        else:
            unassigned.append((idx, song))

    ordered_groups = []
    for sec in sections:
        section_songs = songs_by_section.get(sec["id"], [])
        if section_songs:
            ordered_groups.append((sec["id"], sec["title"], section_songs))

    if unassigned:
        ordered_groups.append((None, None, unassigned))

    return ordered_groups




def format_song_label(song, idx=None):
    if song.get("umh_number"):
        base = f'UMH {song["umh_number"]} {song["title"]}'.strip()
    else:
        base = song.get("title", "")
    if idx is not None:
        return f"{idx+1}. {base}"
    return base


def create_empty_service_order_blocks_from_template(template_sections):
    return [
        {
            "block_id": f"block_{i}",
            "section_id": sec["id"],
            "section_title": sec["title"],
            "source_heading": sec["title"],
            "match_type": "template",
            "score": 100,
            "items": [],
        }
        for i, sec in enumerate(template_sections)
    ]


def pick_default_service_section(template_sections):
    if not template_sections:
        return None

    skip_canonicals = {
        "welcome to mci",
        "welcome",
        "announcements",
    }

    preferred_canonicals = {
        "call to worship",
        "scripture reading",
        "corporate prayer",
        "sermon",
        "response",
        "benediction",
        "tithe offering",
        "tithes offerings",
        "doxology",
    }

    for sec in template_sections:
        canon = canonicalize_section_label(sec["title"])
        if canon in preferred_canonicals:
            return sec

    for sec in template_sections:
        canon = canonicalize_section_label(sec["title"])
        if canon not in skip_canonicals:
            return sec

    return template_sections[0]


def flatten_blocks_to_setlist(service_order_blocks, song_store):
    flat = []
    for block in service_order_blocks:
        for item in block.get("items", []):
            if item.get("type") != "song":
                continue
            song = song_store.get(item.get("song_id"))
            if not song:
                continue
            flat.append({
                "umh_number": song.get("umh_number", ""),
                "title": song.get("title", ""),
                "slides": song.get("slides", []),
                "lyrics_font_size_pt": song.get("lyrics_font_size_pt"),
                "line_spacing": song.get("line_spacing"),
                "override_lyrics_font_size": song.get("override_lyrics_font_size", False),
                "override_line_spacing": song.get("override_line_spacing", False),
                "section_id": block.get("section_id"),
                "song_id": song.get("song_id"),
                "service_block_id": block.get("block_id"),
                "service_block_title": block.get("section_title"),
            })
    return flat


def sync_block_model_from_setlist():
    template_sections = st.session_state.get("template_sections", []) or []
    if not template_sections:
        return

    current_blocks = st.session_state.get("service_order_blocks", []) or []
    song_store = st.session_state.get("song_store", {}) or {}
    setlist = st.session_state.get("setlist", []) or []

    if not current_blocks:
        current_blocks = create_empty_service_order_blocks_from_template(template_sections)

    # preserve non-song items from existing blocks, but rebuild song nesting from setlist
    block_lookup = {}
    new_blocks = []
    for sec in template_sections:
        existing = next((b for b in current_blocks if b.get("section_id") == sec["id"]), None)
        if existing:
            items = [it for it in existing.get("items", []) if it.get("type") != "song"]
            block = {**existing, "section_title": sec["title"], "items": items}
        else:
            block = {
                "block_id": f"block_{len(new_blocks)}",
                "section_id": sec["id"],
                "section_title": sec["title"],
                "source_heading": sec["title"],
                "match_type": "template",
                "score": 100,
                "items": [],
            }
        new_blocks.append(block)
        block_lookup[sec["id"]] = block

    unassigned_items = []
    for idx, song in enumerate(setlist):
        song_id = song.get("song_id") or f"setlist_song_{idx}"
        song_store[song_id] = {
            "song_id": song_id,
            "umh_number": song.get("umh_number", ""),
            "title": song.get("title", ""),
            "slides": song.get("slides", []),
            "lyrics_font_size_pt": song.get("lyrics_font_size_pt"),
            "line_spacing": song.get("line_spacing"),
            "override_lyrics_font_size": song.get("override_lyrics_font_size", False),
            "override_line_spacing": song.get("override_line_spacing", False),
        }
        sec_id = song.get("section_id")
        song_item = {"type": "song", "song_id": song_id}
        if sec_id in block_lookup:
            block_lookup[sec_id]["items"].append(song_item)
        else:
            unassigned_items.append(song_item)

    if unassigned_items:
        new_blocks.append({
            "block_id": "block_unassigned",
            "section_id": None,
            "section_title": "Unassigned Songs",
            "source_heading": "Unassigned Songs",
            "match_type": "derived",
            "score": 100,
            "items": unassigned_items,
        })

    st.session_state["service_order_blocks"] = new_blocks
    st.session_state["song_store"] = song_store


def build_template_service_order_view(setlist=None):
    sync_block_model_from_setlist()
    service_order_blocks = st.session_state.get("service_order_blocks", []) or []
    song_store = st.session_state.get("song_store", {}) or {}
    setlist = st.session_state.get("setlist", []) or []

    index_by_song_id = {}
    for i, song in enumerate(setlist):
        song_id = song.get("song_id")
        if song_id:
            index_by_song_id[song_id] = i

    groups = []
    for block in service_order_blocks:
        items = []
        for item in block.get("items", []):
            if item.get("type") != "song":
                continue
            song = song_store.get(item.get("song_id"))
            if not song:
                continue
            items.append({
                "type": "song",
                "index": index_by_song_id.get(item.get("song_id")),
                "song": song,
            })
        groups.append({
            "section_id": block.get("section_id"),
            "section_title": block.get("section_title"),
            "items": items,
        })
    return groups


def render_service_group_items_markdown(group, selected_index=None, editing_index=None):
    st.markdown(f"**{group['section_title']}**")
    for item in group.get("items", []):
        i = item.get("index")
        song = item.get("song")
        if not song:
            continue
        label = format_song_label(song, i) if i is not None else format_song_label(song)
        prefix = ""
        if selected_index is not None and i == selected_index:
            prefix += "🔹 "
        if editing_index is not None and i == editing_index:
            prefix += "✏️ "
        if prefix:
            st.markdown(f"&nbsp;&nbsp;**{prefix}{label}**", unsafe_allow_html=True)
        else:
            st.markdown(f"&nbsp;&nbsp;{label}", unsafe_allow_html=True)


def render_service_group_items_caption(group):
    st.markdown(f"**{group['section_title']}**")
    for item in group.get("items", []):
        song = item.get("song")
        if not song:
            continue
        st.caption(f"{format_song_label(song)} ({len(song['slides'])} slide(s))")


def reset_editor():
    st.session_state["loaded_song"] = None
    st.session_state["editing_setlist_index"] = None
    st.session_state["editor_umh"] = ""
    st.session_state["editor_title"] = ""
    st.session_state["editor_text"] = ""
    st.session_state["editor_override_lyrics_font_size"] = False
    st.session_state["editor_override_line_spacing"] = False
    st.session_state["editor_lyrics_font_size_pt"] = 32
    st.session_state["editor_line_spacing"] = 1.2
    st.session_state["current_song_preview_images"] = None
    st.session_state["current_song_preview_stats"] = None
    st.session_state["last_editor_text"] = ""
    st.session_state["last_current_song_signature"] = None
    st.session_state["editor_status_message"] = ""
    st.session_state["current_preview_slide"] = 1
    st.session_state["editor_ace_key"] += 1
    st.session_state["selected_song_section_id"] = None


def load_song_preview_if_possible():
    if (
        st.session_state.get("selected_template_name")
        and st.session_state["selected_template_name"] in st.session_state["uploaded_templates"]
        and soffice_available()
    ):
        try:
            template_bytes = st.session_state["uploaded_templates"][
                st.session_state["selected_template_name"]
            ]
            template_ok, _, _ = validate_template_bytes(template_bytes)

            if template_ok:
                current_slides = get_current_slides(st.session_state["editor_text"])
                if current_slides:
                    song_item = build_editor_song_item(current_slides)
                    refresh_current_song_preview(song_item, template_bytes)
                    st.session_state["current_preview_slide"] = 1
                    st.session_state["preview_mode"] = "song"
                    st.session_state["editor_status_message"] = "Song preview loaded."
                else:
                    st.session_state["editor_status_message"] = "Song loaded, but no slides detected for preview."
            else:
                st.session_state["editor_status_message"] = "Song loaded, but selected template is invalid."
        except Exception as e:
            st.session_state["editor_status_message"] = preview_error_message(e)
    else:
        if not st.session_state.get("selected_template_name"):
            st.session_state["editor_status_message"] = "Song loaded. Please select a template to generate preview."
        elif not soffice_available():
            st.session_state["editor_status_message"] = "Song loaded. LibreOffice/soffice is not available."


def load_song_into_editor(match):
    lyrics_raw = str(match.get("Lyrics (Raw)", "")).strip()

    st.session_state["loaded_song"] = {
        "umh_number": str(match.get("UMH Number", "")).strip(),
        "title": str(match.get("Title", "")).strip(),
        "lyrics_raw": lyrics_raw,
    }

    st.session_state["editing_setlist_index"] = None
    st.session_state["editor_umh"] = str(match.get("UMH Number", "")).strip()
    st.session_state["editor_title"] = str(match.get("Title", "")).strip()
    st.session_state["editor_text"] = lyrics_raw

    st.session_state["editor_override_lyrics_font_size"] = False
    st.session_state["editor_override_line_spacing"] = False
    st.session_state["editor_lyrics_font_size_pt"] = 32
    st.session_state["editor_line_spacing"] = 1.2

    sections = st.session_state.get("template_sections", [])
    st.session_state["selected_song_section_id"] = sections[0]["id"] if sections else None

    st.session_state["current_song_preview_images"] = None
    st.session_state["current_song_preview_stats"] = None
    st.session_state["last_editor_text"] = lyrics_raw
    st.session_state["last_current_song_signature"] = None
    st.session_state["editor_status_message"] = ""
    st.session_state["current_preview_slide"] = 1
    st.session_state["preview_mode"] = "song"
    st.session_state["editor_ace_key"] += 1

    load_song_preview_if_possible()


def apply_pending_setlist_load():
    idx = st.session_state.get("pending_setlist_load")
    if idx is None:
        return

    if idx >= len(st.session_state["setlist"]):
        st.session_state["pending_setlist_load"] = None
        return

    item = st.session_state["setlist"][idx]
    lyrics_text = "\n\n".join("\n".join(slide) for slide in item["slides"])

    st.session_state["loaded_song"] = {
        "umh_number": item["umh_number"],
        "title": item["title"],
        "lyrics_raw": lyrics_text,
    }

    st.session_state["editing_setlist_index"] = idx
    st.session_state["editor_umh"] = item["umh_number"]
    st.session_state["editor_title"] = item["title"]
    st.session_state["editor_text"] = lyrics_text

    st.session_state["editor_override_lyrics_font_size"] = item.get(
        "override_lyrics_font_size", False
    )
    st.session_state["editor_override_line_spacing"] = item.get(
        "override_line_spacing", False
    )
    st.session_state["editor_lyrics_font_size_pt"] = item.get("lyrics_font_size_pt", 32) or 32
    st.session_state["editor_line_spacing"] = item.get("line_spacing", 1.2) or 1.2
    st.session_state["selected_song_section_id"] = item.get("section_id")

    if (
        st.session_state["selected_song_section_id"] is None
        and st.session_state.get("template_sections")
    ):
        st.session_state["selected_song_section_id"] = st.session_state["template_sections"][0]["id"]

    st.session_state["current_song_preview_images"] = None
    st.session_state["current_song_preview_stats"] = None
    st.session_state["pending_setlist_load"] = None
    st.session_state["last_editor_text"] = lyrics_text
    st.session_state["last_current_song_signature"] = None
    st.session_state["editor_status_message"] = ""
    st.session_state["current_preview_slide"] = 1
    st.session_state["preview_mode"] = "song"
    st.session_state["editor_ace_key"] += 1

    load_song_preview_if_possible()



# =========================================================
# ORDER OF SERVICE DOCX IMPORT
# =========================================================
UMH_IMPORT_RE = re.compile(
    r"^(?P<hymnal>UMH)\s*(?P<number>\d+)\.?\s+(?P<title>.+?)(?:\s*\((?P<stanzas>.+?)\))?\s*$",
    re.IGNORECASE,
)

DOCX_SECTION_ALIASES = {
    "songs of praise": [
        "songs of praise",
        "hymns of praise",
        "opening song",
        "opening songs",
        "opening hymn",
        "opening hymns",
        "praise songs",
        "praise and worship",
        "worship songs",
    ],
    "gloria patri": ["gloria patri", "glory be to the father"],
    "scripture reading": ["scripture reading", "bible reading", "scripture"],
    "sermon": ["sermon", "message", "sermon title", "the message"],
    "call to worship": ["call to worship"],
    "tithes & offerings": [
        "tithes & offerings",
        "tithes and offerings",
        "offertory",
        "offering",
        "offering prayer",
    ],
    "closing song": [
        "closing song",
        "closing songs",
        "closing hymn",
        "closing hymns",
        "response song",
        "response hymn",
        "sending song",
    ],
    "benediction": ["benediction", "closing blessing"],
}


def normalize_text(s: str) -> str:
    return re.sub(r"\s+", " ", str(s or "").strip()).lower()



def simplify_heading_text(s: str) -> str:
    """
    Normalize headings so variations like:
    - 'CALL TO\u000bWORSHIP'
    - 'Call to Worship (based on Ps 130:5-7)'
    all become:
    - 'call to worship'
    """
    s = str(s or "")

    # Remove vertical tabs from PPT titles
    s = s.replace("\u000b", " ")

    # Remove parenthetical notes
    s = re.sub(r"\([^)]*\)", " ", s)

    # Normalize
    s = s.lower().strip()

    # Normalize symbols
    s = s.replace("&", " and ")

    # Remove punctuation
    s = re.sub(r"[^a-z0-9 ]+", " ", s)

    # Collapse whitespace
    s = re.sub(r"\s+", " ", s).strip()

    return s


def canonicalize_section_label(label: str) -> str:
    simplified = simplify_heading_text(label)
    if not simplified:
        return ""

    for canonical, aliases in DOCX_SECTION_ALIASES.items():
        alias_pool = [canonical] + aliases
        simplified_aliases = [simplify_heading_text(a) for a in alias_pool]
        if simplified in simplified_aliases:
            return canonical
        for alias in simplified_aliases:
            if alias and (simplified == alias or simplified in alias or alias in simplified):
                return canonical

    return simplified



def heading_tokens(s: str) -> set[str]:
    return {tok for tok in simplify_heading_text(s).split() if tok}



def score_section_match(docx_heading: str, template_title: str) -> tuple[int, str]:
    """
    Score similarity between DOCX heading and template section title.
    Returns (score, match_type).
    """
    docx_simple = simplify_heading_text(docx_heading)
    tmpl_simple = simplify_heading_text(template_title)

    if not docx_simple or not tmpl_simple:
        return 0, "none"

    # Exact match after normalization
    if docx_simple == tmpl_simple:
        return 100, "exact"

    docx_tokens = set(docx_simple.split())
    tmpl_tokens = set(tmpl_simple.split())

    # Same tokens but different order/format
    if docx_tokens == tmpl_tokens:
        return 96, "alias"

    # Substring containment
    if docx_simple in tmpl_simple or tmpl_simple in docx_simple:
        return 90, "contains"

    # Partial token overlap
    overlap = docx_tokens & tmpl_tokens
    if overlap:
        score = 60 + min(20, 10 * len(overlap))
        return score, "token-overlap"

    return 0, "none"


def read_docx_lines(docx_file) -> list[str]:
    doc = Document(docx_file)
    lines = []
    for para in doc.paragraphs:
        text = para.text.replace(" ", " ").strip()
        if text:
            lines.append(text)
    return lines



def is_umh_song_line(line: str) -> bool:
    return UMH_IMPORT_RE.match(line.strip()) is not None



def parse_umh_song_line(line: str) -> dict:
    m = UMH_IMPORT_RE.match(line.strip())
    if not m:
        return {
            "umh_number": "",
            "title": line.strip(),
            "stanzas": "",
            "raw": line.strip(),
        }

    return {
        "umh_number": (m.group("number") or "").strip(),
        "title": (m.group("title") or "").strip(),
        "stanzas": (m.group("stanzas") or "").strip(),
        "raw": line.strip(),
    }



def match_template_section_from_heading(heading: str, template_sections: list[dict]) -> dict | None:
    """
    Robust heading-to-template matcher.

    Handles cases like:
    - "Call to Worship (based on Ps 130:5-7)"
    - "CALL TO\u000bWORSHIP"
    """
    if not template_sections or not heading:
        return None

    heading = apply_docx_heading_alias(heading)
    heading_simple = simplify_heading_text(heading)
    
    if not heading_simple:
        return None

    best = None

    for sec in template_sections:
        template_simple = simplify_heading_text(sec["title"])

        score = 0
        match_type = "none"

        if heading_simple == template_simple:
            score = 100
            match_type = "exact"

        elif heading_simple.startswith(template_simple) or template_simple.startswith(heading_simple):
            score = 98
            match_type = "alias"

        elif heading_simple in template_simple or template_simple in heading_simple:
            score = 95
            match_type = "contains"

        else:
            heading_tokens = set(heading_simple.split())
            template_tokens = set(template_simple.split())
            overlap = heading_tokens & template_tokens

            if heading_tokens and template_tokens and overlap:
                ratio = len(overlap) / max(1, min(len(heading_tokens), len(template_tokens)))
                if ratio >= 0.75:
                    score = 90
                    match_type = "alias"
                elif ratio >= 0.5:
                    score = 75
                    match_type = "token-overlap"

        candidate = {
            "section_id": sec["id"],
            "section_title": sec["title"],
            "score": score,
            "match_type": match_type,
            "docx_heading": heading,
        }

        if best is None or candidate["score"] > best["score"]:
            best = candidate

    if best and best["score"] >= 90:
        return best

    return None

def section_id_from_heading(heading: str, template_sections: list[dict]) -> str | None:
    match = match_template_section_from_heading(heading, template_sections)
    return match["section_id"] if match else None



def build_song_item_from_row(row, section_id=None):
    lyrics_raw = str(row.get("Lyrics (Raw)", "")).strip()
    slides = split_slides_manual(lyrics_raw)

    return {
        "umh_number": str(row.get("UMH Number", "")).strip(),
        "title": str(row.get("Title", "")).strip(),
        "slides": slides,
        "lyrics_font_size_pt": None,
        "line_spacing": None,
        "override_lyrics_font_size": False,
        "override_line_spacing": False,
        "section_id": section_id,
    }



def should_treat_docx_line_as_anchor(line: str) -> bool:
    stripped = str(line or "").strip()
    if not stripped:
        return False
    if is_umh_song_line(stripped):
        return False
    if stripped.startswith("+") or stripped.startswith("#"):
        return True

    # Allow plain section titles from the DOCX to act as range anchors.
    # This lets the importer map everything between major headings like
    # Call to Worship -> Scripture Reading -> Corporate Prayer.
    word_count = len(stripped.split())
    has_sentence_punctuation = any(ch in stripped for ch in [":", ";", ".", "?", "!"])
    if word_count <= 6 and not has_sentence_punctuation:
        return True
    return False



def import_service_order_from_docx(docx_file, template_sections):
    """
    Template-driven block import.

    Returns:
      service_order_blocks: list[dict]
      song_store: dict[str, dict]
      missing_songs: list[dict]
      section_mapping_rows: list[dict]
    """
    lines = read_docx_lines(docx_file)

    service_order_blocks = create_empty_service_order_blocks_from_template(template_sections)
    block_by_section_id = {block["section_id"]: block for block in service_order_blocks}
    song_store = {}
    missing_songs = []
    section_mapping_rows = []

    current_block = None
    song_seq = 0

    def pick_starting_block():
        default_sec = pick_default_service_section(template_sections)
        if default_sec:
            return block_by_section_id.get(default_sec["id"])
        return service_order_blocks[0] if service_order_blocks else None

    current_block = pick_starting_block()

    for raw_line in lines:
        stripped = raw_line.strip()
        if not stripped:
            continue

        is_prefixed_heading = stripped.startswith(("+", "#"))
        heading_text = stripped[1:].strip() if is_prefixed_heading else stripped

        if not is_prefixed_heading:
            heading_match = match_template_section_from_heading(heading_text, template_sections)
            if heading_match and heading_match["match_type"] in {"exact", "alias"}:
                matched_block = block_by_section_id.get(heading_match["section_id"])
                if matched_block is not None:
                    current_block = matched_block
                    section_mapping_rows.append({
                        "docx_heading": heading_text,
                        "mapped_section_id": heading_match["section_id"],
                        "mapped_section_title": heading_match["section_title"],
                        "match_type": heading_match["match_type"],
                        "score": heading_match["score"],
                        "source": "plain",
                    })
                    continue

        if is_prefixed_heading:
            if current_block is not None:
                current_block["items"].append({"type": "minor_heading", "text": heading_text})
            continue

        if is_umh_song_line(stripped):
            parsed = parse_umh_song_line(stripped)
            row = find_row_by_umh(parsed["umh_number"])
            if not row:
                missing_songs.append(parsed)
                continue

            song_seq += 1
            song_id = f"song_{song_seq}"
            song_store[song_id] = {
                "song_id": song_id,
                "umh_number": str(row.get("UMH Number", "")).strip(),
                "title": str(row.get("Title", "")).strip(),
                "slides": split_slides_manual(str(row.get("Lyrics (Raw)", "")).strip()),
                "lyrics_font_size_pt": None,
                "line_spacing": None,
                "override_lyrics_font_size": False,
                "override_line_spacing": False,
            }
            if current_block is not None:
                current_block["items"].append({"type": "song", "song_id": song_id})
            continue

        if current_block is not None:
            current_block["items"].append({"type": "text", "text": stripped})

    return service_order_blocks, song_store, missing_songs, section_mapping_rows


# =========================================================
# PPT BUILD HELPERS
# =========================================================
def add_song_block_to_prs(prs, song, first_layout, rest_layout):
    start_idx = len(prs.slides)

    umh_number = str(song["umh_number"]).strip()
    title = str(song["title"]).strip()
    slides = song["slides"]

    lyrics_font_size_pt = song.get("lyrics_font_size_pt")
    line_spacing = song.get("line_spacing")

    full_title = f"UMH {umh_number} {title}".strip() if umh_number else title

    for i, slide_lines in enumerate(slides):
        lyrics_text = "\n".join(slide_lines)

        if i == 0:
            slide = prs.slides.add_slide(first_layout)
            set_shape_text(slide.shapes.title, full_title)
            set_shape_text(
                get_body_placeholder(slide),
                lyrics_text,
                font_size_pt=lyrics_font_size_pt,
                line_spacing=line_spacing,
            )
        else:
            slide = prs.slides.add_slide(rest_layout)
            set_shape_text(
                get_body_placeholder(slide),
                lyrics_text,
                font_size_pt=lyrics_font_size_pt,
                line_spacing=line_spacing,
            )

    end_idx = len(prs.slides) - 1
    return start_idx, end_idx


def add_section_song_block_to_prs(prs, section_song_pairs, first_layout, rest_layout):
    start_idx = len(prs.slides)
    for _, song in section_song_pairs:
        add_song_block_to_prs(prs, song, first_layout, rest_layout)
    end_idx = len(prs.slides) - 1
    return start_idx, end_idx


def create_combined_ppt(setlist, template_bytes: bytes):
    sync_block_model_from_setlist()
    service_order_blocks = st.session_state.get("service_order_blocks", []) or []
    song_store = st.session_state.get("song_store", {}) or {}

    prs = open_presentation_from_bytes(template_bytes)

    first_layout = get_layout_by_name(prs, FIRST_LAYOUT_NAME)
    rest_layout = get_layout_by_name(prs, REST_LAYOUT_NAME)

    if first_layout is None or rest_layout is None:
        raise ValueError("Template layouts not found.")

    output_mode = st.session_state.get("service_output_mode", "full")

    ordered_song_pairs = []
    for block in service_order_blocks:
        for item in block.get("items", []):
            if item.get("type") == "song":
                song = song_store.get(item.get("song_id"))
                if song:
                    ordered_song_pairs.append((item["song_id"], song))

    if output_mode == "songs":
        delete_all_slides(prs)
        for _, song in ordered_song_pairs:
            add_song_block_to_prs(prs, song, first_layout, rest_layout)

        output = BytesIO()
        prs.save(output)
        output.seek(0)
        return output

    if not st.session_state.get("preserve_template_slides", True):
        delete_all_slides(prs)
        for _, song in ordered_song_pairs:
            add_song_block_to_prs(prs, song, first_layout, rest_layout)

        output = BytesIO()
        prs.save(output)
        output.seek(0)
        return output

    hidden_ids = set(st.session_state.get("hidden_section_ids", []))
    sections = st.session_state.get("template_sections", [])

    slides_to_delete = []
    for sec in sections:
        if sec["id"] in hidden_ids:
            slides_to_delete.extend(sec.get("content_slide_indices", []))

    for idx in sorted(set(slides_to_delete), reverse=True):
        delete_slide_by_index(prs, idx)

    for block in service_order_blocks:
        sec_title = block.get("section_title")
        section_song_pairs = []
        for item in block.get("items", []):
            if item.get("type") != "song":
                continue
            song = song_store.get(item.get("song_id"))
            if song:
                section_song_pairs.append((item["song_id"], song))

        if not section_song_pairs:
            continue

        if block.get("section_id") is None:
            add_section_song_block_to_prs(prs, section_song_pairs, first_layout, rest_layout)
            continue

        target_idx = find_section_insert_index(prs, sec_title)
        start_idx, end_idx = add_section_song_block_to_prs(
            prs,
            section_song_pairs,
            first_layout,
            rest_layout,
        )
        move_slide_block(prs, start_idx, end_idx, target_idx)

    output = BytesIO()
    prs.save(output)
    output.seek(0)
    return output


def create_single_song_ppt(song_item, template_bytes: bytes):
    prs = open_presentation_from_bytes(template_bytes)

    first_layout = get_layout_by_name(prs, FIRST_LAYOUT_NAME)
    rest_layout = get_layout_by_name(prs, REST_LAYOUT_NAME)

    if first_layout is None or rest_layout is None:
        raise ValueError("Template layouts not found.")

    delete_all_slides(prs)
    add_song_block_to_prs(prs, song_item, first_layout, rest_layout)

    output = BytesIO()
    prs.save(output)
    output.seek(0)
    return output


# =========================================================
# PREVIEW HELPERS
# =========================================================
def pptx_to_preview_images(pptx_bytes: BytesIO):
    with tempfile.TemporaryDirectory() as tmpdir:
        pptx_path = os.path.join(tmpdir, "preview.pptx")

        with open(pptx_path, "wb") as f:
            f.write(pptx_bytes.getvalue())

        cmd = [
            SOFFICE_PATH,
            "--headless",
            "--convert-to", "pdf",
            "--outdir", tmpdir,
            pptx_path,
        ]

        result = subprocess.run(cmd, capture_output=True, text=True)

        stderr_text = (result.stderr or "").strip().lower()

        if (
            result.returncode != 0
            or "source file could not be loaded" in stderr_text
            or "error:" in stderr_text
        ):
            raise RuntimeError(
                "Preview conversion failed.\n"
                f"Return code: {result.returncode}\n"
                f"stdout: {result.stdout}\n"
                f"stderr: {result.stderr}"
            )

        pdf_files = [f for f in os.listdir(tmpdir) if f.lower().endswith(".pdf")]
        if not pdf_files:
            raise FileNotFoundError(
                "Preview PDF was not created.\n"
                f"Files in temp dir: {os.listdir(tmpdir)}\n"
                f"stdout: {result.stdout}\n"
                f"stderr: {result.stderr}"
            )

        pdf_path = os.path.join(tmpdir, pdf_files[0])

        try:
            doc = fitz.open(pdf_path)
            repaired_bytes = doc.tobytes(garbage=3, clean=True, deflate=True)
            doc.close()
            doc = fitz.open(stream=repaired_bytes, filetype="pdf")
        except Exception as e:
            raise RuntimeError(f"Unable to open preview PDF in PyMuPDF: {e}")

        images = []
        try:
            for page in doc:
                pix = page.get_pixmap(dpi=60, alpha=False)
                img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)

                buffer = io.BytesIO()
                img.save(buffer, format="JPEG", quality=45, optimize=True)
                images.append(buffer.getvalue())

                buffer.close()
                del pix
                del img
                del buffer
        finally:
            doc.close()

        return images


def preview_error_message(exc: Exception) -> str:
    msg = str(exc).strip()

    if "source file could not be loaded" in msg:
        return (
            "LibreOffice could not open the generated PowerPoint preview file. "
            "This usually means the combined PPTX structure is invalid for LibreOffice."
        )

    if "Preview conversion failed" in msg:
        return "Preview conversion failed."

    if "Preview PDF was not created" in msg:
        return "Preview PDF was not created."

    if "No common ancestor in structure tree" in msg:
        return (
            "Preview generation failed because the generated PDF structure is malformed."
        )

    return f"Preview generation failed: {msg}"


def render_scrollable_images(images, height=760, active_slide=None):
    if not images:
        st.info("Preview will appear here.")
        return

    container_id = f"preview-scroll-container-{len(images)}-{active_slide}"
    active_slide_js = "null" if active_slide is None else str(active_slide)

    html = f"""
    <div id="{container_id}" style="
        height: {height}px;
        overflow-y: auto;
        border: 1px solid #ddd;
        padding: 12px;
        border-radius: 8px;
        background: #fafafa;
        box-sizing: border-box;
        scroll-behavior: smooth;
    ">
    """

    for i, img_bytes in enumerate(images, start=1):
        b64 = base64.b64encode(img_bytes).decode("utf-8")
        border = "3px solid #2563eb" if active_slide == i else "1px solid #ccc"
        badge = " ← current" if active_slide == i else ""

        html += f"""
        <div id="slide-{i}" style="margin-bottom: 24px;">
            <div style="font-weight: 600; margin-bottom: 8px;">Slide {i}{badge}</div>
            <img
                src="data:image/jpeg;base64,{b64}"
                style="width: 100%; border: {border}; display: block;"
            />
        </div>
        """

    html += "</div>"

    html += f"""
    <script>
    const container = document.getElementById("{container_id}");
    const activeSlide = {active_slide_js};

    function scrollToActiveSlide() {{
        if (!container || activeSlide === null) return;
        const target = document.getElementById("slide-" + activeSlide);
        if (!target) return;
        container.scrollTop = target.offsetTop - 12;
    }}

    setTimeout(() => {{
        scrollToActiveSlide();
    }}, 200);
    </script>
    """

    st.components.v1.html(html, height=height, scrolling=False)


# =========================================================
# AUTO PREVIEW HELPERS
# =========================================================
def refresh_current_song_preview(song_item, template_bytes):
    st.session_state["current_song_preview_images"] = None
    st.session_state["current_song_preview_stats"] = None

    ppt_data = create_single_song_ppt(song_item, template_bytes)
    preview_images = pptx_to_preview_images(ppt_data)

    st.session_state["current_song_preview_images"] = preview_images
    st.session_state["current_song_preview_stats"] = preview_stats(preview_images)
    st.session_state["last_current_song_signature"] = build_current_song_signature(
        song_item,
        st.session_state.get("selected_template_name"),
    )


def blank_separator_added(old_text: str, new_text: str) -> bool:
    def valid_blank_positions(lines):
        positions = []
        for i, line in enumerate(lines):
            if line.strip() == "":
                for j in range(i + 1, len(lines)):
                    if lines[j].strip() != "":
                        positions.append(i)
                        break
        return positions

    old_lines = old_text.splitlines()
    new_lines = new_text.splitlines()
    return len(valid_blank_positions(new_lines)) > len(valid_blank_positions(old_lines))


def detect_new_slide_target_line(old_text: str, new_text: str):
    old_lines = old_text.splitlines()
    new_lines = new_text.splitlines()

    min_len = min(len(old_lines), len(new_lines))
    changed_idx = None

    for i in range(min_len):
        if old_lines[i] != new_lines[i]:
            changed_idx = i
            break

    if changed_idx is None:
        if len(new_lines) > len(old_lines):
            changed_idx = len(old_lines)
        else:
            return None

    for j in range(changed_idx, len(new_lines)):
        if new_lines[j].strip() != "":
            return j

    return None


def get_first_new_blank_separator_index(old_text: str, new_text: str):
    old_blank_positions = [i for i, line in enumerate(old_text.splitlines()) if line.strip() == ""]
    new_blank_positions = [i for i, line in enumerate(new_text.splitlines()) if line.strip() == ""]

    for pos in new_blank_positions:
        if pos not in old_blank_positions:
            return pos
    return None


def get_slide_number_from_line_index(text: str, line_index: int, auto_split: bool, lines_per_slide: int):
    if line_index is None:
        return None

    lines = text.splitlines()

    if auto_split:
        current_verse_indexes = []
        line_to_slide = {}
        slide_num = 1

        for idx, raw_line in enumerate(lines):
            stripped = raw_line.strip()

            if stripped == "":
                if current_verse_indexes:
                    for j in range(0, len(current_verse_indexes), lines_per_slide):
                        chunk = current_verse_indexes[j:j + lines_per_slide]
                        for original_idx in chunk:
                            line_to_slide[original_idx] = slide_num
                        slide_num += 1
                    current_verse_indexes = []
            else:
                current_verse_indexes.append(idx)

        if current_verse_indexes:
            for j in range(0, len(current_verse_indexes), lines_per_slide):
                chunk = current_verse_indexes[j:j + lines_per_slide]
                for original_idx in chunk:
                    line_to_slide[original_idx] = slide_num
                slide_num += 1

        return line_to_slide.get(line_index)

    slide_num = 1
    in_slide = False

    for idx, raw_line in enumerate(lines):
        stripped = raw_line.strip()

        if stripped == "":
            if in_slide:
                slide_num += 1
                in_slide = False
        else:
            in_slide = True
            if idx == line_index:
                return slide_num

    return None


# =========================================================
# SERVICE PREVIEW POSITION HELPERS
# =========================================================
def get_service_song_start_slides(setlist, template_bytes: bytes):
    starts = [None] * len(setlist)

    if st.session_state.get("service_output_mode") == "songs" or not st.session_state.get("preserve_template_slides", True):
        slide_counter = 1
        ordered = get_ordered_songs_for_output(setlist)
        for original_idx, song in ordered:
            starts[original_idx] = slide_counter
            slide_counter += len(song["slides"])
        return starts

    sections = st.session_state.get("template_sections", [])
    hidden_ids = set(st.session_state.get("hidden_section_ids", []))
    groups = group_songs_by_section_order(setlist)

    songs_by_sec = {sec_id: song_pairs for sec_id, _, song_pairs in groups if sec_id is not None}
    unassigned = next((song_pairs for sec_id, _, song_pairs in groups if sec_id is None), [])

    slide_counter = 1

    for sec in sections:
        slide_counter += 1
        if sec["id"] not in hidden_ids:
            slide_counter += len(sec.get("content_slide_indices", []))

        for original_idx, song in songs_by_sec.get(sec["id"], []):
            starts[original_idx] = slide_counter
            slide_counter += len(song["slides"])

    for original_idx, song in unassigned:
        starts[original_idx] = slide_counter
        slide_counter += len(song["slides"])

    return starts


def refresh_service_preview(setlist, template_bytes):
    sync_block_model_from_setlist()
    st.session_state["service_preview_images"] = None
    st.session_state["service_preview_stats"] = None
    st.session_state["service_preview_error"] = None
    st.session_state["ppt_data"] = None

    ppt_data = create_combined_ppt(setlist, template_bytes)
    st.session_state["ppt_data"] = ppt_data
    st.session_state["service_song_start_slides"] = get_service_song_start_slides(
        setlist, template_bytes
    )

    try:
        preview_images = pptx_to_preview_images(ppt_data)

        if not preview_images:
            raise RuntimeError("No preview images generated from service PPT")

        st.session_state["service_preview_images"] = preview_images
        st.session_state["service_preview_stats"] = preview_stats(preview_images)

    except Exception as e:
        st.session_state["service_preview_images"] = None
        st.session_state["service_preview_stats"] = None
        st.session_state["service_preview_error"] = preview_error_message(e)


def selected_template_info():
    template_name = st.session_state.get("selected_template_name")
    if not template_name:
        return None, False, [], []

    uploaded = st.session_state["uploaded_templates"]
    if template_name not in uploaded:
        return None, False, [], []

    template_bytes = uploaded[template_name]
    ok, errors, warnings = validate_template_bytes(template_bytes)
    return template_bytes, ok, errors, warnings


# =========================================================
# PRE-RUN ACTIONS
# =========================================================
apply_pending_setlist_load()

if st.session_state.get("reset_editor_pending"):
    reset_editor()
    st.session_state["reset_editor_pending"] = False

selected_template_bytes, selected_template_ok, selected_template_errors, selected_template_warnings = selected_template_info()

if selected_template_bytes and selected_template_ok:
    try:
        current_sections = parse_template_sections(selected_template_bytes)
        old_sections = st.session_state.get("template_sections", [])

        old_hidden_by_title = {
            sec["title"]: (sec["id"] in st.session_state.get("hidden_section_ids", []))
            for sec in old_sections
        }

        st.session_state["template_sections"] = current_sections

        new_hidden_ids = []
        for sec in current_sections:
            if old_hidden_by_title.get(sec["title"], False):
                new_hidden_ids.append(sec["id"])
        st.session_state["hidden_section_ids"] = new_hidden_ids

    except Exception:
        st.session_state["template_sections"] = []
        st.session_state["hidden_section_ids"] = []
else:
    st.session_state["template_sections"] = []
    st.session_state["hidden_section_ids"] = []

if st.session_state.get("template_sections"):
    valid_section_ids = {sec["id"] for sec in st.session_state["template_sections"]}
    if st.session_state.get("selected_song_section_id") not in valid_section_ids:
        st.session_state["selected_song_section_id"] = st.session_state["template_sections"][0]["id"]
else:
    st.session_state["selected_song_section_id"] = None

if st.session_state.get("template_sections") and not st.session_state.get("service_order_blocks"):
    st.session_state["service_order_blocks"] = create_empty_service_order_blocks_from_template(
        st.session_state["template_sections"]
    )


# =========================================================
# SIDEBAR
# =========================================================
with st.sidebar:
    sync_block_model_from_setlist()
    st.markdown("### Service Order")

    setlist = st.session_state["setlist"]

    if not setlist:
        st.caption("No songs added yet.")
    else:
        selected_index = st.session_state.get("setlist_selected_index", 0)
        editing_index = st.session_state.get("editing_setlist_index")

        selected_index = max(0, min(selected_index, len(setlist) - 1))
        st.session_state["setlist_selected_index"] = selected_index

        for group in build_template_service_order_view(setlist):
            render_service_group_items_markdown(
                group,
                selected_index=selected_index,
                editing_index=editing_index,
            )

    st.divider()

    st.markdown("### Loaded in Editor")

    loaded_umh = st.session_state.get("editor_umh", "").strip()
    loaded_title = st.session_state.get("editor_title", "").strip()
    loaded_idx = st.session_state.get("editing_setlist_index")
    loaded_section_title = get_section_title_by_id(st.session_state.get("selected_song_section_id"))

    if loaded_title:
        section_line = f"\n\nSection: {loaded_section_title}" if loaded_section_title else ""

        if loaded_idx is not None:
            if loaded_umh:
                st.info(f"Editing setlist item #{loaded_idx + 1}\n\nUMH {loaded_umh} {loaded_title}{section_line}")
            else:
                st.info(f"Editing setlist item #{loaded_idx + 1}\n\n{loaded_title}{section_line}")
        else:
            if loaded_umh:
                st.info(f"New / repository song\n\nUMH {loaded_umh} {loaded_title}{section_line}")
            else:
                st.info(f"New / repository song\n\n{loaded_title}{section_line}")
    else:
        st.caption("No song loaded in editor.")

    st.divider()

    st.header("Controls")

    with st.expander("1. Template", expanded=False):
        uploaded_templates = st.file_uploader(
            "Upload template(s)",
            type=["pptx"],
            accept_multiple_files=True,
            key="template_uploader",
        )

        if uploaded_templates:
            for file in uploaded_templates:
                st.session_state["uploaded_templates"][file.name] = file.getvalue()

        template_names = list(st.session_state["uploaded_templates"].keys())

        if template_names:
            default_index = 0
            if st.session_state["selected_template_name"] in template_names:
                default_index = template_names.index(st.session_state["selected_template_name"])

            chosen_template = st.selectbox(
                "Select template",
                template_names,
                index=default_index,
            )

            if chosen_template != st.session_state.get("selected_template_name"):
                st.session_state["selected_template_name"] = chosen_template
                clear_service_outputs()
                st.session_state["current_song_preview_images"] = None
                st.session_state["current_song_preview_stats"] = None
                st.session_state["last_current_song_signature"] = None
                st.rerun()
            else:
                st.session_state["selected_template_name"] = chosen_template

            selected_template_bytes, selected_template_ok, selected_template_errors, selected_template_warnings = selected_template_info()

            if selected_template_ok:
                st.success("Template valid")
            else:
                st.error("Template invalid")
                for err in selected_template_errors:
                    st.write(f"- {err}")

            if selected_template_warnings:
                for warn in selected_template_warnings:
                    st.warning(warn)

            if st.button("Remove selected template", use_container_width=True):
                del st.session_state["uploaded_templates"][st.session_state["selected_template_name"]]
                st.session_state["selected_template_name"] = None
                st.session_state["template_sections"] = []
                st.session_state["hidden_section_ids"] = []
                st.session_state["selected_song_section_id"] = None
                clear_service_outputs()
                st.rerun()

            st.divider()
            st.checkbox(
                "Keep existing slides in template",
                key="preserve_template_slides",
            )

            st.radio(
                "Service output mode",
                options=["full", "songs"],
                format_func=lambda x: (
                    "Full deck (insert section-by-section)"
                    if x == "full"
                    else "Songs only (for checking)"
                ),
                key="service_output_mode",
            )

            if selected_template_ok and st.session_state["template_sections"]:
                st.divider()
                st.markdown("#### Hide Template Section Contents")

                section_options = [sec["id"] for sec in st.session_state["template_sections"]]
                section_name_map = {
                    sec["id"]: sec["title"] for sec in st.session_state["template_sections"]
                }

                st.multiselect(
                    "Hide contents under these section headers",
                    options=section_options,
                    default=st.session_state.get("hidden_section_ids", []),
                    format_func=lambda sid: section_name_map[sid],
                    key="hidden_section_ids",
                )

                for sec in st.session_state["template_sections"]:
                    hidden = sec["id"] in st.session_state["hidden_section_ids"]
                    status = "content hidden" if hidden else "content shown"
                    st.caption(
                        f'{sec["title"]}: {len(sec["content_slide_indices"])} content slide(s) · divider kept · {status}'
                    )
            elif selected_template_ok:
                st.info(
                    f'No "{DIVIDER_LAYOUT_NAME}" slides detected. '
                    "Section tools will not be available."
                )
        else:
            st.info("Upload at least one template.")

        if not soffice_available():
            st.warning("LibreOffice/soffice is not available.")

    with st.expander("2. Load Song", expanded=False):
        if st.button("Start New Song", use_container_width=True):
            st.session_state["reset_editor_pending"] = True
            st.rerun()

        load_mode = st.radio("Find hymn by", ["UMH Number", "Title"], horizontal=True)

        if load_mode == "UMH Number":
            umh_number_input = st.text_input("UMH Number", placeholder="e.g. 57")
            if st.button("Load by Number", use_container_width=True):
                if umh_number_input.strip():
                    match = find_row_by_umh(umh_number_input)
                    if match:
                        load_song_into_editor(match)
                        st.success("Hymn loaded.")
                        st.rerun()
                    else:
                        st.error("Hymn not found.")
        else:
            keyword = st.text_input("Search title", placeholder="e.g. thousand tongues")
            matches = search_titles(keyword) if keyword.strip() else []

            if matches:
                options = [
                    f'UMH {row.get("UMH Number","")} - {row.get("Title","")}'
                    for row in matches
                ]
                selected = st.selectbox("Select hymn", options)

                if st.button("Load by Title", use_container_width=True):
                    chosen_index = options.index(selected)
                    load_song_into_editor(matches[chosen_index])
                    st.success("Hymn loaded.")
                    st.rerun()
            elif keyword.strip():
                st.info("No matching titles found.")

    with st.expander("3. Service Order", expanded=True):
        setlist = st.session_state["setlist"]

        if not setlist:
            st.info("No songs added yet.")
        else:
            st.markdown("#### Service Order (by DOCX block, using template section headers)")
            for group in build_template_service_order_view(setlist):
                if not group["items"]:
                    continue
                render_service_group_items_caption(group)

            st.divider()
            st.markdown("#### Song Selection")
            labels = []
            for i, song in enumerate(setlist):
                section_title = get_section_title_by_id(song.get("section_id"))
                section_suffix = f" [{section_title}]" if section_title else ""

                if song["umh_number"]:
                    labels.append(f'{i+1}. UMH {song["umh_number"]} {song["title"]}{section_suffix} ({len(song["slides"])})')
                else:
                    labels.append(f'{i+1}. {song["title"]}{section_suffix} ({len(song["slides"])})')

            st.session_state["setlist_selected_index"] = min(
                st.session_state.get("setlist_selected_index", 0),
                len(labels) - 1,
            )

            pending_index = st.session_state.pop("pending_setlist_selectbox_index", None)

            if pending_index is None:
                pending_index = st.session_state.get("setlist_selectbox_sidebar", 0)

            try:
                pending_index = int(pending_index)
            except Exception:
                pending_index = 0

            pending_index = max(0, min(pending_index, len(labels) - 1))

            if "setlist_selectbox_sidebar" not in st.session_state:
                st.session_state["setlist_selectbox_sidebar"] = pending_index
            st.session_state["setlist_selected_index"] = pending_index

            previous_selected_index = st.session_state["setlist_selected_index"]

            selected_index = st.selectbox(
                "Selected song",
                options=list(range(len(labels))),
                format_func=lambda i: labels[i],
                key="setlist_selectbox_sidebar",
            )
            
            try:
                selected_index = int(selected_index)
            except Exception:
                selected_index = 0
            
            selected_index = max(0, min(selected_index, len(labels) - 1))
            st.session_state["setlist_selected_index"] = selected_index
            
            if 0 <= selected_index < len(setlist):
                st.session_state["selected_song_id"] = setlist[selected_index].get("song_id")

            if (
                selected_index != previous_selected_index
                and st.session_state.get("preview_mode") == "service"
            ):
                starts = st.session_state.get("service_song_start_slides", [])
                st.session_state["current_preview_slide"] = (
                    starts[selected_index] if selected_index < len(starts) and starts[selected_index] is not None else 1
                )
                st.rerun()

            action_cols = st.columns(4)

            with action_cols[0]:
                if st.button("✏️", use_container_width=True, help="Edit selected song"):
                    st.session_state["pending_setlist_load"] = selected_index
                    st.session_state["preview_mode"] = "song"
                    st.session_state["current_song_preview_images"] = None
                    st.session_state["current_song_preview_stats"] = None
                    st.session_state["last_current_song_signature"] = None
                    st.rerun()

            with action_cols[1]:
                if st.button("⬆️", use_container_width=True, help="Move selected song up") and selected_index > 0:
                    setlist[selected_index - 1], setlist[selected_index] = (
                        setlist[selected_index],
                        setlist[selected_index - 1],
                    )

                    new_index = selected_index - 1
                    st.session_state["setlist_selected_index"] = new_index
                    st.session_state["pending_setlist_selectbox_index"] = new_index

                    editing_index = st.session_state.get("editing_setlist_index")
                    if editing_index == selected_index:
                        st.session_state["editing_setlist_index"] = new_index
                    elif editing_index == new_index:
                        st.session_state["editing_setlist_index"] = selected_index

                    pending = st.session_state.get("pending_setlist_load")
                    if pending == selected_index:
                        st.session_state["pending_setlist_load"] = new_index
                    elif pending == new_index:
                        st.session_state["pending_setlist_load"] = selected_index

                    clear_service_outputs()
                    st.rerun()

            with action_cols[2]:
                if st.button("⬇️", use_container_width=True, help="Move selected song down") and selected_index < len(setlist) - 1:
                    setlist[selected_index + 1], setlist[selected_index] = (
                        setlist[selected_index],
                        setlist[selected_index + 1],
                    )

                    new_index = selected_index + 1
                    st.session_state["setlist_selected_index"] = new_index
                    st.session_state["pending_setlist_selectbox_index"] = new_index

                    editing_index = st.session_state.get("editing_setlist_index")
                    if editing_index == selected_index:
                        st.session_state["editing_setlist_index"] = new_index
                    elif editing_index == new_index:
                        st.session_state["editing_setlist_index"] = selected_index

                    pending = st.session_state.get("pending_setlist_load")
                    if pending == selected_index:
                        st.session_state["pending_setlist_load"] = new_index
                    elif pending == new_index:
                        st.session_state["pending_setlist_load"] = selected_index

                    clear_service_outputs()
                    st.rerun()

            with action_cols[3]:
                if st.button("🗑️", use_container_width=True, help="Delete selected song"):
                    setlist.pop(selected_index)

                    new_index = min(selected_index, len(setlist) - 1) if setlist else 0
                    st.session_state["setlist_selected_index"] = new_index
                    st.session_state["pending_setlist_selectbox_index"] = new_index

                    if st.session_state.get("editing_setlist_index") == selected_index:
                        st.session_state["reset_editor_pending"] = True
                    elif (
                        st.session_state.get("editing_setlist_index") is not None
                        and st.session_state["editing_setlist_index"] > selected_index
                    ):
                        st.session_state["editing_setlist_index"] -= 1

                    pending = st.session_state.get("pending_setlist_load")
                    if pending == selected_index:
                        st.session_state["pending_setlist_load"] = None
                    elif pending is not None and pending > selected_index:
                        st.session_state["pending_setlist_load"] = pending - 1

                    clear_service_outputs()
                    st.rerun()

            if st.button("Clear Setlist", use_container_width=True, type="secondary"):
                st.session_state["setlist"] = []
                st.session_state["service_order_blocks"] = []
                st.session_state["song_store"] = {}
                st.session_state["last_docx_section_mapping"] = []
                st.session_state["editing_setlist_index"] = None
                st.session_state["pending_setlist_load"] = None
                st.session_state["setlist_selected_index"] = 0
                st.session_state["pending_setlist_selectbox_index"] = None
                st.session_state["preview_mode"] = "song"
                st.session_state["current_song_preview_images"] = None
                st.session_state["current_song_preview_stats"] = None
                clear_service_outputs()
                st.rerun()


    with st.expander("4. Import Order of Service", expanded=False):
        service_docx_file = st.file_uploader(
            "Upload Order of Service (.docx)",
            type=["docx"],
            key="service_docx_importer",
        )

        replace_existing_setlist = st.checkbox(
            "Replace current setlist",
            value=False,
            key="replace_setlist_from_docx",
        )
        st.caption("Template sections remain the source of truth. DOCX is parsed into ordered blocks, and only major matched headings start a new template section block.")

        if st.button("Import Songs from DOCX", use_container_width=True):
            if service_docx_file is None:
                st.warning("Please upload a .docx file first.")
            else:
                try:
                    service_order_blocks, song_store, missing_songs, section_mapping_rows = import_service_order_from_docx(
                        service_docx_file,
                        st.session_state.get("template_sections", []),
                    )

                    st.session_state["service_order_blocks"] = service_order_blocks
                    st.session_state["song_store"] = song_store
                    st.session_state["setlist"] = flatten_blocks_to_setlist(service_order_blocks, song_store)
                    st.session_state["last_docx_section_mapping"] = section_mapping_rows

                    clear_service_outputs()

                    st.session_state["editing_setlist_index"] = None
                    st.session_state["pending_setlist_load"] = None
                    st.session_state["setlist_selected_index"] = 0
                    st.session_state["pending_setlist_selectbox_index"] = 0

                    st.success(f"Imported {len(st.session_state["setlist"])} song(s) into the template-driven service order.")

                    if section_mapping_rows:
                        st.caption("Section mapping:")
                        for row in section_mapping_rows:
                            mapped = row.get("mapped_section_title") or "No template section matched"
                            st.write(
                                f"{row['source']} {row['docx_heading']} → {mapped}"
                                f" ({row['match_type']}, score {row['score']})"
                            )

                    if missing_songs:
                        st.warning(
                            f"{len(missing_songs)} song(s) were found in the DOCX but not found in Google Sheets."
                        )
                        for song in missing_songs:
                            st.write(f"- UMH {song['umh_number']} {song['title']}")

                    st.rerun()

                except Exception as e:
                    st.exception(e)



# =========================================================
# MAIN LAYOUT
# =========================================================
main_left, main_right = st.columns([1.15, 1], vertical_alignment="top")

with main_left:
    st.subheader("Song Editor")

    edit_idx = st.session_state.get("editing_setlist_index")
    if edit_idx is not None:
        st.info(f"Editing setlist item #{edit_idx + 1}")

    meta_col1, meta_col2 = st.columns([1, 3])
    with meta_col1:
        st.text_input("UMH", key="editor_umh")
    with meta_col2:
        st.text_input("Title", key="editor_title")

    if st.session_state.get("template_sections"):
        section_options = [sec["id"] for sec in st.session_state["template_sections"]]
        section_name_map = {sec["id"]: sec["title"] for sec in st.session_state["template_sections"]}

        current_section = st.session_state.get("selected_song_section_id")
        if current_section not in section_options:
            current_section = section_options[0]
            st.session_state["selected_song_section_id"] = current_section

        st.selectbox(
            "Place new song under section",
            options=section_options,
            format_func=lambda sid: section_name_map[sid],
            key="selected_song_section_id",
        )
    else:
        st.caption("No template sections detected.")

    st.markdown("#### Slide Splitting")
    split_col1, split_col2 = st.columns([3, 1])

    with split_col1:
        st.checkbox("Auto split by lines per slide", key="auto_split_by_lines")
        st.slider("Lines per slide", min_value=1, max_value=8, key="lines_per_slide")
        st.checkbox(
            "Refresh preview only when a new slide break is detected",
            key="refresh_on_new_line",
        )

    with split_col2:
        st.write("")
        refresh_song_preview_clicked = st.button("Refresh Song Preview", use_container_width=True)

    old_text = st.session_state.get("last_editor_text", "")
    editor_text = st_ace(
        value=st.session_state.get("editor_text", ""),
        language="text",
        theme="textmate",
        keybinding="vscode",
        font_size=16,
        tab_size=2,
        wrap=True,
        show_gutter=False,
        auto_update=True,
        readonly=False,
        height=420,
        key=f"editor_ace_{st.session_state['editor_ace_key']}",
    )

    if editor_text is None:
        editor_text = st.session_state.get("editor_text", "")
    else:
        st.session_state["editor_text"] = editor_text

    if editor_text != old_text:
        clear_service_outputs()

    current_slides = get_current_slides(editor_text)
    song_item = build_editor_song_item(current_slides)

    if st.session_state["auto_split_by_lines"]:
        st.caption(
            f"{len(current_slides)} slide(s) "
            f"({st.session_state['lines_per_slide']} lines per slide, blank lines kept as verse separators)"
        )
    else:
        st.caption(f"{len(current_slides)} slide(s) (manual mode: blank lines separate slides)")

    if refresh_song_preview_clicked:
        if selected_template_bytes is None:
            st.error("Please upload and select a template first.")
        elif not selected_template_ok:
            st.error("Selected template is invalid.")
        elif not soffice_available():
            st.error("LibreOffice/soffice is not available.")
        elif not current_slides:
            st.error("No slides to preview.")
        else:
            try:
                st.session_state["preview_mode"] = "song"
                st.session_state["current_song_preview_images"] = None
                st.session_state["current_song_preview_stats"] = None
                refresh_current_song_preview(song_item, selected_template_bytes)
                st.session_state["editor_status_message"] = "Song preview refreshed."
                st.rerun()
            except Exception as e:
                st.error(preview_error_message(e))

    text_changed = editor_text != old_text
    trigger_refresh = False
    current_split_settings = (
        st.session_state["auto_split_by_lines"],
        st.session_state["lines_per_slide"],
    )

    if st.session_state["auto_split_by_lines"]:
        trigger_refresh = get_current_slides(old_text) != get_current_slides(editor_text)
    else:
        trigger_refresh = blank_separator_added(old_text, editor_text)

    if text_changed and trigger_refresh:
        if not st.session_state["auto_split_by_lines"]:
            blank_idx = get_first_new_blank_separator_index(old_text, editor_text)
            if blank_idx is not None:
                lines = editor_text.splitlines()
                target_line_index = None
                for i in range(blank_idx + 1, len(lines)):
                    if lines[i].strip() != "":
                        target_line_index = i
                        break

                detected_slide = get_slide_number_from_line_index(
                    editor_text,
                    target_line_index,
                    auto_split=False,
                    lines_per_slide=st.session_state["lines_per_slide"],
                )
                if detected_slide is not None:
                    st.session_state["current_preview_slide"] = detected_slide
        else:
            target_line_index = detect_new_slide_target_line(old_text, editor_text)
            detected_slide = get_slide_number_from_line_index(
                editor_text,
                target_line_index,
                auto_split=True,
                lines_per_slide=st.session_state["lines_per_slide"],
            )
            if detected_slide is not None:
                st.session_state["current_preview_slide"] = detected_slide

    st.session_state["last_editor_text"] = editor_text
    st.session_state["last_split_settings"] = current_split_settings

    if st.session_state["editor_status_message"]:
        st.caption(st.session_state["editor_status_message"])

    st.markdown("#### Song Formatting")

    fmt_col1, fmt_col2 = st.columns(2)

    with fmt_col1:
        st.checkbox("Override lyrics font size", key="editor_override_lyrics_font_size")
        if st.session_state["editor_override_lyrics_font_size"]:
            st.slider(
                "Lyrics font size (pt)",
                min_value=12,
                max_value=60,
                key="editor_lyrics_font_size_pt",
                on_change=lambda: st.session_state.update({"last_current_song_signature": None}),
            )
        else:
            st.caption("Using template lyrics size")

    with fmt_col2:
        st.checkbox("Override line spacing", key="editor_override_line_spacing")
        if st.session_state["editor_override_line_spacing"]:
            st.slider(
                "Line spacing",
                min_value=0.8,
                max_value=2.0,
                step=0.1,
                key="editor_line_spacing",
                on_change=lambda: st.session_state.update({"last_current_song_signature": None}),
            )
        else:
            st.caption("Using template line spacing")

    st.markdown("#### Add / Update")
    allow_duplicates = st.checkbox("Allow duplicate songs in setlist", value=False)

    action_col1, action_col2 = st.columns(2)
    with action_col1:
        add_or_update = st.button(
            "Update Song" if edit_idx is not None else "Add to Setlist",
            use_container_width=True,
        )
    with action_col2:
        clear_editor = st.button("Clear Editor", use_container_width=True)

    if add_or_update:
        if not current_slides:
            st.error("No slides to add.")
        else:
            item = build_editor_song_item(current_slides)
            edit_idx = st.session_state.get("editing_setlist_index")

            if edit_idx is None:
                duplicate_index = next(
                    (
                        i for i, s in enumerate(st.session_state["setlist"])
                        if s["umh_number"] == item["umh_number"]
                        and s["title"] == item["title"]
                        and s["slides"] == item["slides"]
                        and s.get("section_id") == item.get("section_id")
                    ),
                    None,
                )

                if duplicate_index is not None and not allow_duplicates:
                    st.warning(f"This song is already in the setlist as item #{duplicate_index + 1}.")
                else:
                    st.session_state["setlist"].append(item)
                    clear_service_outputs()
                    st.success(
                        f'Added: {"UMH " + item["umh_number"] + " " if item["umh_number"] else ""}{item["title"]}'
                    )
                    st.session_state["reset_editor_pending"] = True
                    st.rerun()
            else:
                old_item = st.session_state["setlist"][edit_idx]
                for meta_key in ["import_uid", "service_block_id", "service_block_order", "service_block_title"]:
                    if meta_key in old_item:
                        item[meta_key] = old_item[meta_key]
                if not item.get("section_id") and old_item.get("section_id"):
                    item["section_id"] = old_item.get("section_id")
                st.session_state["setlist"][edit_idx] = item
                clear_service_outputs()

                if (
                    selected_template_bytes is not None
                    and selected_template_ok
                    and soffice_available()
                    and current_slides
                ):
                    try:
                        st.session_state["current_song_preview_images"] = None
                        st.session_state["current_song_preview_stats"] = None
                        refresh_current_song_preview(item, selected_template_bytes)
                        st.session_state["current_preview_slide"] = 1
                        st.session_state["preview_mode"] = "song"
                        st.session_state["editor_status_message"] = "Song updated and preview refreshed."
                    except Exception as e:
                        st.session_state["editor_status_message"] = preview_error_message(e)

                st.success(
                    f'Updated: {"UMH " + item["umh_number"] + " " if item["umh_number"] else ""}{item["title"]}'
                )
                st.rerun()

    if clear_editor:
        st.session_state["reset_editor_pending"] = True
        st.rerun()


with main_right:
    st.subheader("Preview")

    resource_stats = get_runtime_resource_stats()
    with st.expander("Runtime resource usage", expanded=False):
        st.write(f'App RAM (RSS): {format_bytes(resource_stats["process_rss"])}')
        st.write(f'App virtual memory: {format_bytes(resource_stats["process_vms"])}')
        st.write(
            f'System RAM: {format_bytes(resource_stats["system_used_ram"])} used / '
            f'{format_bytes(resource_stats["system_total_ram"])} total'
        )
        st.write(
            f'System RAM available now: {format_bytes(resource_stats["system_available_ram"])}'
        )
        st.write(
            f'Disk: {format_bytes(resource_stats["disk_used"])} used / '
            f'{format_bytes(resource_stats["disk_total"])} total'
        )
        st.write(f'Disk free now: {format_bytes(resource_stats["disk_free"])}')

    preview_mode = st.radio(
        "Preview Mode",
        ["Song", "Service"],
        index=0 if st.session_state["preview_mode"] == "song" else 1,
        horizontal=True,
    )
    st.session_state["preview_mode"] = preview_mode.lower()

    if st.session_state["preview_mode"] == "song":
        st.session_state["service_preview_images"] = None
        st.session_state["service_preview_stats"] = None
    else:
        st.session_state["current_song_preview_images"] = None
        st.session_state["current_song_preview_stats"] = None

    preview_images = None

    if st.session_state["preview_mode"] == "song":
        st.caption("Current song preview")

        if selected_template_bytes is None:
            st.warning("Please upload and select a template to see preview.")
        elif not selected_template_ok:
            st.error("Selected template is invalid.")
        elif not soffice_available():
            st.warning("LibreOffice/soffice is required for preview.")

        preview_images = st.session_state.get("current_song_preview_images")

        song_stats = st.session_state.get("current_song_preview_stats")
        if song_stats:
            st.caption(
                f'Preview: {song_stats["count"]} slide(s) · '
                f'total {format_bytes(song_stats["total_bytes"])} · '
                f'avg {format_bytes(song_stats["avg_bytes"])} · '
                f'max {format_bytes(song_stats["max_bytes"])}'
            )

    else:
        if st.session_state.get("service_output_mode") == "songs":
            st.caption("Output mode: songs only (for checking)")
        else:
            st.caption("Output mode: full deck with section-by-section insertion")

        can_generate = (
            selected_template_bytes is not None
            and selected_template_ok
            and soffice_available()
            and len(st.session_state["setlist"]) > 0
        )

        if can_generate:
            refresh_service_now = st.button(
                "Refresh Service Preview",
                use_container_width=True
            )

            if refresh_service_now:
                try:
                    refresh_service_preview(
                        st.session_state["setlist"],
                        selected_template_bytes,
                    )
                except Exception as e:
                    st.error(f"PowerPoint generation failed: {e}")

            starts = st.session_state.get("service_song_start_slides", [])
            selected_index = st.session_state.get("setlist_selected_index", 0)

            if starts and 0 <= selected_index < len(starts) and starts[selected_index] is not None:
                st.session_state["current_preview_slide"] = starts[selected_index]
            else:
                st.session_state["current_preview_slide"] = 1

            if st.session_state.get("service_preview_error"):
                st.warning(
                    "Preview could not be generated, but the PowerPoint file is available for download.\n\n"
                    + st.session_state["service_preview_error"]
                )

            preview_images = st.session_state.get("service_preview_images")

            service_stats = st.session_state.get("service_preview_stats")
            if service_stats:
                st.caption(
                    f'Preview: {service_stats["count"]} slide(s) · '
                    f'total {format_bytes(service_stats["total_bytes"])} · '
                    f'avg {format_bytes(service_stats["avg_bytes"])} · '
                    f'max {format_bytes(service_stats["max_bytes"])}'
                )

            if st.session_state.get("ppt_data") is not None:
                st.download_button(
                    label="Download Service PowerPoint",
                    data=st.session_state["ppt_data"].getvalue(),
                    file_name="service_deck.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    use_container_width=True,
                )

        else:
            if selected_template_bytes is None:
                st.info("Please upload and select a template first.")
            elif not selected_template_ok:
                st.info("Selected template is invalid.")
            elif not soffice_available():
                st.info("LibreOffice/soffice is not available.")
            elif not st.session_state["setlist"]:
                st.info("Add songs to the setlist to view the service preview.")

    if preview_images:
        render_scrollable_images(
            preview_images,
            height=700,
            active_slide=st.session_state.get("current_preview_slide"),
        )
    else:
        st.info("Preview will appear here.")
